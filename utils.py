import os
import fitz  # PyMuPDF
import mammoth
import sys
import subprocess
import re
import pdftables_api
from docx import Document
from config import apiKey
from pdf2docx import Converter
from docxcompose.composer import Composer
from pdf2image import convert_from_path

import pandas as pd
import ppt2pdf
from docx import Document
from htmldocx import HtmlToDocx

from pptx import Presentation
from pptx.util import Inches
import PyPDF2



def count_pdf_pages(pdf_file_path):
    try:
        pdf_document = fitz.open(pdf_file_path)
        num_pages = pdf_document.page_count
        return num_pages
    except Exception as e:
        print(f"An error occurred: {str(e)}")
        return None


def getFilename(path):
    file = path.split("/")[-1]
    file = file.split(".")[0]
    return file
    

def pdf2doc(path, outputPath):
    pdf_file = path
    filename = getFilename(path)
    # docx_file = f"{outputPath}/{filename}.docx"

    num_pages = count_pdf_pages(pdf_file)

    for i in range(num_pages):
        docx_file = f'{outputPath}/{filename}{i}.docx'

        # convert pdf to docx
        cv = Converter(pdf_file)
        cv.convert(docx_file, start=i, end= i + 1)      # all pages by default
        cv.close()
    
    master = Document(f"{outputPath}/{filename}{0}.docx")
    composer = Composer(master)
    os.remove(f"{outputPath}/{filename}{0}.docx")
        
    for i in range(1, num_pages):
        doc = Document(f"{outputPath}/{filename}{i}.docx")
        composer.append(doc)
        os.remove(f"{outputPath}/{filename}{i}.docx")
        
    docx_file = f"{outputPath}/{filename}.docx"
    composer.save(docx_file)
    
    return docx_file

def pdf2ppt(path, outputPath):
    os.system(f"pdf2pptx {path}")
    output = path.replace("pdf", "pptx")
    return output

def ppt2pdf(path, outputPath):
    os.system(f"ppt2pdf file {path}")
    outputPath = path.replace("pptx", "pdf")
    outputPath = path.replace("ppt", "pdf")
    return outputPath
    
    


# pip install git+https://github.com/pdftables/python-pdftables-api.git
# pip setup.py install 
def pdf2csv(path, outputPath):
    c = pdftables_api.Client(apiKey)
    filename = getFilename(path)
    c.xlsx(path, f'{outputPath}/{filename}.xlsx') 
    return f'{outputPath}/{filename}.xlsx'

def pdf2html(path, outputPath):
    filepath = pdf2doc(path, outputPath)
    
    with open(filepath, "rb") as docx_file:
        result = mammoth.convert_to_html(docx_file)
        html = result.value 
        messages = result.messages
    
    os.remove(filepath)
    filename = getFilename(path)
    out = f"{outputPath}/{filename}.html"
    f = open(out,"w")
    f.write(html)
    f.close()
    return out


def pdf2JPEG(path, outputPath):
    pages = convert_from_path(path)
    filename = getFilename(path)
    out = f"{outputPath}/{filename}.jpg"
    for page in pages:
        page.save(out, 'JPEG') 
    return out



def convert_to(folder, source, timeout=None):
    args = [libreoffice_exec(), '--headless', '--convert-to', 'pdf', '--outdir', folder, source]

    process = subprocess.run(args, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=timeout)
    filename = re.search('-> (.*?) using filter', process.stdout.decode())

    if filename is None:
        raise LibreOfficeError(process.stdout.decode())
    else:
        return filename.group(1)
    


def libreoffice_exec():
    # TODO: Provide support for more platforms
    if sys.platform == 'darwin':
        return '/Applications/LibreOffice.app/Contents/MacOS/soffice'
    return 'libreoffice'


class LibreOfficeError(Exception):
    def __init__(self, output):
        self.output = output
        

def docx2pdf(path, outputPath):
    filename = getFilename(path)
    out = f"{outputPath}/{filename}.pdf"
    convert_to(outputPath, path)
    
    return out


def csv2pdf(path, outputPath):
    
    filename = getFilename(path)
    df = pd.read_excel(path)
    html_string = df.to_html()
    
    document = Document()
    new_parser = HtmlToDocx()
    
    html = html_string
    new_parser.add_html_to_document(html, document)

    out = f"output/{filename}.docx"
    # do more stuff to document
    document.save(out)
    
    ret = docx2pdf(out, outputPath)
    os.remove(out)
    
    return ret


def html2pdf(path, outputPath):
    file = open(path, 'r').read()
    filename = getFilename(path)
    
    out = f"{outputPath}/{filename}.docx"
    
    document = Document()
    new_parser = HtmlToDocx()
    # do stuff to document

    new_parser.add_html_to_document(file, document)

    # do more stuff to document
    document.save(out)
    
    ret = docx2pdf(out, outputPath)
    os.remove(out)
    
    return ret







def extract_text_from_pdf(pdf_file):
    text = ""
    listt = []
    with open(pdf_file, 'rb') as pdf:
        pdf_reader = PyPDF2.PdfReader(pdf)
        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            listt.append(page.extract_text())
    return listt


def create_presentation(strings, output_file):
    # Create a new presentation
    X = Presentation()
    Layout = X.slide_layouts[0]
    
    first_slide = X.slides.add_slide(Layout)
    text = strings[0]
    title_text = text.split('\n')
    title = first_slide.shapes.title
    title.text = title_text[0]

    first_slide.shapes.title.text = title_text[0]
    first_slide.placeholders[1].text = '\n'.join(title_text[1:])

    for text in strings[1:]:   
        title_text = text.split('\n')     
        Second_Layout = X.slide_layouts[5]
        second_slide = X.slides.add_slide(Second_Layout)
        second_slide.shapes.title.text = title_text[0]

        textbox = second_slide.shapes.add_textbox(Inches(3), Inches(1.5),Inches(3), Inches(1)) 
        textframe = textbox.text_frame
        paragraph = textframe.add_paragraph()
        paragraph.text = '\n'.join(title_text[1:])

    # Save the presentation to the specified file
    X.save(output_file)
    
    
def pdf2PPT(path, outputPath):
    filename = getFilename(path)
    out = f"{outputPath}/{filename}.pptx"
    
    texts = extract_text_from_pdf(path)
    create_presentation(texts, out)
    
    return True

    
    
    